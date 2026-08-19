from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

INPUT = 'Concept-to-Prototype-Deck.pptx'
OUTPUT = 'Concept-to-Prototype-Deck-styled.pptx'

prs = Presentation(INPUT)

# Theme colors
BG = RGBColor(11, 57, 84)       # deep blue
TITLE = RGBColor(255, 255, 255)  # white
BODY = RGBColor(230, 230, 230)   # light gray
ACCENT = RGBColor(255, 179, 0)   # accent amber (not used by default)

for slide in prs.slides:
    # set solid background color
    try:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BG
    except Exception:
        pass

    # style title if present
    try:
        title = slide.shapes.title
        if title and title.has_text_frame:
            for p in title.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(30)
                    run.font.bold = True
                    run.font.color.rgb = TITLE
    except Exception:
        pass

    # style other text frames (body)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        # skip title placeholder
        try:
            if slide.shapes.title and shape == slide.shapes.title:
                continue
        except Exception:
            pass

        tf = shape.text_frame
        for p in tf.paragraphs:
            # set first paragraph as heading-like if it looks like a subtitle (heuristic)
            for run in p.runs:
                run.font.size = Pt(18)
                run.font.color.rgb = BODY

# Save styled presentation
prs.save(OUTPUT)
print('Saved styled PPTX:', OUTPUT)
