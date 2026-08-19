import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt

INPUT = 'Concept-to-Prototype-Deck.md'
OUTPUT = 'Concept-to-Prototype-Deck.pptx'

if not os.path.exists(INPUT):
    print(f"Input file not found: {INPUT}")
    raise SystemExit(1)

with open(INPUT, 'r', encoding='utf-8') as f:
    text = f.read()

# Split on markdown slide separator lines '---' on their own
slides_raw = re.split(r'\n-{3,}\n', text)

prs = Presentation()
# Remove default first blank slide
if prs.slides:
    # keep the blank slide layout removed by adding content later
    pass

for block in slides_raw:
    block = block.strip()
    if not block:
        continue
    lines = [ln.rstrip() for ln in block.splitlines() if ln.strip()]
    if not lines:
        continue
    title = lines[0]
    body_lines = lines[1:]

    # Choose a Title and Content layout when possible
    try:
        slide_layout = prs.slide_layouts[1]
    except Exception:
        slide_layout = prs.slide_layouts[0]

    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = title

    # Prepare content text frame
    # Try to use placeholder 1
    text_items = []
    for ln in body_lines:
        ln = ln.strip()
        # treat list items starting with '-' as bullets
        m = re.match(r'^[-*]\s+(.*)', ln)
        if m:
            text_items.append(('bullet', m.group(1).strip()))
        else:
            text_items.append(('text', ln))

    # Populate placeholder if exists
    try:
        placeholder = slide.placeholders[1]
        tf = placeholder.text_frame
        tf.clear()
        first = True
        for kind, content in text_items:
            if first:
                p = tf.paragraphs[0]
                p.text = content
                p.level = 0
                p.font.size = Pt(18)
                first = False
            else:
                p = tf.add_paragraph()
                p.text = content
                p.level = 0
                p.font.size = Pt(18)
    except Exception:
        # fallback: add a textbox
        left = Inches(1)
        top = Inches(1.8)
        width = Inches(8)
        height = Inches(4.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        for kind, content in text_items:
            p = tf.add_paragraph()
            p.text = content
            p.font.size = Pt(18)

prs.save(OUTPUT)
print(f'Saved: {OUTPUT}')
